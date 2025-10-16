import streamlit as st
import time
import os
import json
from datetime import datetime
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests

# Import your PPTGenerator class here
class PPTGenerator:
    def __init__(self, api_key=None, pexel_key=None):
        self.api_key = api_key
        self.pexel_key = pexel_key
        
        if not self.api_key:
            raise ValueError('Gemini API key is not available')

        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel('gemini-2.5-pro')
        self.presentation = Presentation()
        
        self.last_api_call = 0
        self.min_delay = 30
        self.api_call_count = 0

    def _wait_for_rate_limit(self):
        current_time = time.time()
        time_since_last_call = current_time - self.last_api_call
        
        if time_since_last_call < self.min_delay:
            wait_time = self.min_delay - time_since_last_call
            st.info(f"Rate limiting: Waiting {wait_time:.1f}s before next API call...")
            time.sleep(wait_time)
        
        self.last_api_call = time.time()
        self.api_call_count += 1

    def _retry_api_call(self, func, *args, max_retries=3, **kwargs):
        for attempt in range(max_retries):
            try:
                self._wait_for_rate_limit()
                return func(*args, **kwargs)
            except Exception as e:
                error_msg = str(e)
                
                if "429" in error_msg or "quota" in error_msg.lower():
                    if attempt < max_retries - 1:
                        wait_time = 45 * (attempt + 1)
                        st.warning(f"Rate limit hit. Waiting {wait_time}s before retry {attempt + 2}/{max_retries}...")
                        time.sleep(wait_time)
                        continue
                    else:
                        st.error(f"Max retries reached. Skipping this API call.")
                        return None
                
                elif "503" in error_msg:
                    if attempt < max_retries - 1:
                        wait_time = 5 * (attempt + 1)
                        st.warning(f"Service unavailable (503). Retrying in {wait_time}s...")
                        time.sleep(wait_time)
                        continue
                    else:
                        st.error(f"Service unavailable after {max_retries} attempts.")
                        return None
                else:
                    st.error(f"Error: {e}")
                    return None
        
        return None

    def generate_content_outlines(self, topic, num_slides=5):
        prompt = f"""
        Create a detailed outline for a PowerPoint presentation on "{topic}" with {num_slides} slides.
        Return the response as a JSON array with the following structure:
        [
          {{
            "title": "Slide Title",
            "content": "Main content points as bullet points",
            "slide_type": "title | content | image | conclusion"
          }}
        ]
        Make sure the content is engaging, informative, and well-structured.
        The response must be a valid JSON array.
        """

        def _generate():
            response = self.model.generate_content(prompt)
            return response.text.strip()

        try:
            content = self._retry_api_call(_generate)
            
            if content is None:
                raise ValueError("Failed to generate content after retries")

            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].strip()

            if not content.startswith('['):
                st.warning("Gemini did not return JSON. Using fallback content.")
                return self._get_fallback_content(topic)

            return json.loads(content)

        except Exception as e:
            st.error(f"Error generating content: {e}")
            return self._get_fallback_content(topic)

    def _get_fallback_content(self, topic):
        return [
            {"title": f"{topic}", "content": "Introduction and Overview", "slide_type": "title"},
            {"title": "Background", "content": "Key background information", "slide_type": "content"},
            {"title": "Main Points", "content": "Important details and examples", "slide_type": "content"},
            {"title": "Analysis", "content": "Critical analysis and insights", "slide_type": "content"},
            {"title": "Future Outlook", "content": "Trends and predictions", "slide_type": "content"},
            {"title": "Conclusion", "content": "Summary and takeaways", "slide_type": "conclusion"}
        ]

    def generative_image_description(self, slide_content):
        content_lower = str(slide_content).lower()
        
        keyword_map = {
            'introduction': 'professional business presentation',
            'conclusion': 'success celebration team',
            'future': 'futuristic technology innovation',
            'challenge': 'business strategy planning',
            'benefit': 'growth success chart',
            'process': 'workflow diagram business',
            'team': 'diverse team collaboration',
            'ai': 'artificial intelligence technology',
            'movie': 'film production cinema',
            'production': 'movie set film making',
            'technology': 'modern technology innovation'
        }
        
        for keyword, image_desc in keyword_map.items():
            if keyword in content_lower:
                st.info(f"Using cached image description: {image_desc}")
                return image_desc
        
        prompt = f"""
        Based on this slide content, suggest a relevant image description that would enhance the presentation
        {slide_content}
        Return only a brief, descriptive phrase suitable for image search (max 7 words)
        """
        
        def _generate():
            response = self.model.generate_content(prompt)
            return response.text.strip()
        
        try:
            result = self._retry_api_call(_generate)
            return result if result else 'professional presentation visual'
        except Exception as e:
            st.error(f"Error generating image description: {e}")
            return 'professional presentation visual'

    def download_images(self, query, save_path='temp_image.jpg'):
        try:
            url = 'https://api.pexels.com/v1/search'
            header = {'Authorization': self.pexel_key}
            params = {'query': query, 'per_page': 1, 'orientation': 'landscape'}

            response = requests.get(url, headers=header, params=params)
            response.raise_for_status()

            data = response.json()
            if not data.get('photos'):
                raise ValueError('No photo found')

            image_url = data['photos'][0]['src']['original']
            image_response = requests.get(image_url)
            image_response.raise_for_status()

            with open(save_path, 'wb') as f:
                f.write(image_response.content)

            return save_path

        except Exception as e:
            st.warning(f"Error downloading image: {e}")
            return None

    def create_title_slide(self, title, subtitle):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(30)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def create_content_slide(self, title, content, include_image=False):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(30)
        title_shape.text_frame.paragraphs[0].font.bold = True

        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        if isinstance(content, list):
            content_text = '\n'.join(str(item) for item in content)
        else:
            content_text = str(content)
        
        p = text_frame.paragraphs[0]
        p.text = content_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        if include_image:
            try:
                image_desc = self.generative_image_description(content_text)
                if image_desc:
                    image_path = self.download_images(image_desc)
                    if image_path and os.path.exists(image_path):
                        slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
                        os.remove(image_path)
            except Exception as e:
                st.warning(f"Error adding image to slide: {e}")

        return slide

    def create_image_slide(self, title, content, image_query):
        slide_layout = self.presentation.slide_layouts[8]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(2.5))
        content_frame = content_box.text_frame
        
        if isinstance(content, list):
            content_text = '\n'.join(str(item) for item in content)
        else:
            content_text = str(content)
        
        content_frame.text = content_text
        content_frame.word_wrap = True

        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
            paragraph.alignment = PP_ALIGN.CENTER

        try:
            image_path = self.download_images(image_query)
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(3.25), Inches(4), width=Inches(3.5), height=Inches(2.5))
                os.remove(image_path)
        except Exception as e:
            st.warning(f"Error adding image to slide: {e}")

        return slide

    def generate_ppt(self, topic, num_slides=6, output_file='presentation.pptx', author='Pranab'):
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        status_text.text(f"Starting PPT generation for: {topic}")
        st.info(f"Total slides: {num_slides}")
        st.info(f"Estimated time: ~{num_slides * 30 // 60} minutes (due to rate limits)")
        
        content_outlines = self.generate_content_outlines(topic, num_slides)
        total_slides = len(content_outlines)

        for i, slide_data in enumerate(content_outlines):
            title = slide_data.get('title', f"Slide {i+1}")
            content = slide_data.get('content', "")
            slide_type = slide_data.get('slide_type', 'content')

            status_text.text(f"📄 Generating slide {i+1}/{total_slides}: {title}")
            progress_bar.progress((i + 1) / total_slides)

            if i == 0 or slide_type == 'title':
                self.create_title_slide(title, f'Created by {author}')
            elif slide_type == 'content':
                self.create_content_slide(title, content, include_image=True)
            elif slide_type == 'image':
                content_str = '\n'.join(content) if isinstance(content, list) else str(content)
                img_query = self.generative_image_description(content_str)
                self.create_image_slide(title, content, img_query)
            else:
                self.create_content_slide(title, content, include_image=False)

        self.presentation.save(output_file)
        status_text.text(f"PPT generated successfully!")
        
        
        return output_file


# Streamlit App
def main():
    st.set_page_config(
        page_title="AI PPT Generator",
        layout="wide"
    )
    
    # Custom CSS
    st.markdown("""
        <style>
        .main {
            padding: 2rem;
        }
        .stButton>button {
            width: 100%;
            background-color: #4CAF50;
            color: white;
            font-size: 18px;
            padding: 12px;
            border-radius: 8px;
        }
        .stTextInput>div>div>input {
            font-size: 16px;
        }
        </style>
    """, unsafe_allow_html=True)
    
    # Header
    st.title("AI PowerPoint Generator")
    st.markdown("### Generate professional presentations powered by Google Gemini AI")
    st.divider()
    
    # Sidebar for API keys
    with st.sidebar:
        st.header("Configuration")
        
        gemini_api_key = st.text_input(
            "Gemini API Key",
            type="password",
            help="Get your API key from https://makersuite.google.com/app/apikey"
        )
        
        pexel_api_key = st.text_input(
            "Pexels API Key",
            type="password",
            help="Get your API key from https://www.pexels.com/api/"
        )
        
        st.divider()
        
        author_name = st.text_input(
            "Author Name",
            value="Pranab",
            help="Your name will appear on the title slide"
        )
        
        st.divider()
        st.info("**Tip:** Free tier allows 2 API calls per minute. Generation may take a few minutes.")
    
    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        topic = st.text_input(
            "Presentation Topic",
            placeholder="e.g., AI in Movie Industry, Climate Change, etc.",
            help="Enter the main topic for your presentation"
        )
    
    with col2:
        num_slides = st.number_input(
            "Number of Slides",
            min_value=3,
            max_value=15,
            value=6,
            help="Choose between 3-15 slides"
        )
    
    # Advanced options
    with st.expander("Advanced Options"):
        col3, col4 = st.columns(2)
        
        with col3:
            output_filename = st.text_input(
                "Output Filename",
                value="presentation.pptx",
                help="Name of the generated PowerPoint file"
            )
        
        with col4:
            include_images = st.checkbox(
                "Include Images",
                value=True,
                help="Add relevant images to slides"
            )
    
    st.divider()
    
    # Generate button
    if st.button("Generate Presentation", type="primary"):
        
        # Validation
        if not topic:
            st.error("Please enter a presentation topic!")
            return
        
        if not gemini_api_key:
            st.error("Please provide your Gemini API key in the sidebar!")
            return
        
        if not pexel_api_key:
            st.error("Please provide your Pexels API key in the sidebar!")
            return
        
        # Generate presentation
        try:
            with st.spinner("Creating your presentation..."):
                generator = PPTGenerator(api_key=gemini_api_key, pexel_key=pexel_api_key)
                
                output_file = generator.generate_ppt(
                    topic=topic,
                    num_slides=num_slides,
                    output_file=output_filename,
                    author=author_name
                )
                
                # Download button
                with open(output_file, "rb") as file:
                    st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name=output_filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
                st.balloons()
                
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
            st.exception(e)
    
    # Footer
    st.divider()
    st.markdown("""
        <div style='text-align: center; color: gray;'>
            <p>Made with using Streamlit and Google Gemini AI</p>
            <p>by <a href='https://github.com/decodephi' target='_blank'>Pranab</a></p>
        </div>
    """, unsafe_allow_html=True)


if __name__ == "__main__":
    main()